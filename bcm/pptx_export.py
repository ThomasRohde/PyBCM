import math
from typing import List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from .models import LayoutModel
from .layout import process_layout
from .settings import Settings

def hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def pixels_to_inches(pixels: float) -> float:
    """Convert pixels to inches (assuming 96 DPI)."""
    return pixels / 96.0

def calculate_font_size(root_size: int, level: int, is_leaf: bool) -> int:
    """Calculate font size based on level and node type."""
    # Calculate base size for this level - decrease by 4 points per level
    base_size = root_size - (level * 4)
    
    if is_leaf:
        # Leaf nodes get a slightly smaller size than their parent level
        return max(base_size - 2, 8)
    else:
        # Non-leaf nodes use the level's base size
        return max(base_size, 10)

def add_node_to_slide(slide, node: LayoutModel, settings: Settings, level: int = 0):
    """Add a node and its children to the PowerPoint slide."""
    # Convert coordinates and dimensions to inches
    left = pixels_to_inches(node.x)
    top = pixels_to_inches(node.y)
    width = pixels_to_inches(node.width)
    height = pixels_to_inches(node.height)
    
    # Determine node color based on level and whether it has children
    if not node.children:
        color = settings.get("color_leaf")
    else:
        color = settings.get(f"color_{min(level, 6)}")
    
    # Convert hex color to RGB
    rgb_color = hex_to_rgb(color)
    
    # Add shape
    shape = slide.shapes.add_shape(
        1,  # Rectangle shape
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )
    
    # Set shape fill color
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb_color)
    
    # Set shape line color
    shape.line.color.rgb = RGBColor(51, 51, 51)  # #333333
    shape.line.width = Pt(1)
    
    # Calculate font size
    font_size = calculate_font_size(settings.get("root_font_size"), level, not node.children)
    
    # Add text
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = True
    # Use top alignment for nodes with children, middle for leaf nodes
    text_frame.vertical_anchor = MSO_ANCHOR.TOP if node.children else MSO_ANCHOR.MIDDLE
    
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.text = node.name
    
    # Set font properties
    run = paragraph.runs[0]
    run.font.size = Pt(font_size)
    run.font.name = 'Arial'
    run.font.color.rgb = RGBColor(0, 0, 0)
    
    # Recursively add child nodes
    if node.children:
        for child in node.children:
            add_node_to_slide(slide, child, settings, level + 1)

def export_to_pptx(model: LayoutModel, settings: Settings) -> Presentation:
    """Export the capability model to PowerPoint format."""
    # Process layout
    processed_model = process_layout(model, settings)
    
    # Create presentation
    prs = Presentation()
    
    # Add slide
    slide_layout = prs.slide_layouts[6]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Calculate dimensions with padding
    padding = settings.get("padding", 20)
    width = math.ceil(processed_model.width + 2 * padding)
    height = math.ceil(processed_model.height + 2 * padding)
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(pixels_to_inches(width))
    prs.slide_height = Inches(pixels_to_inches(height))
    
    # Add nodes recursively
    add_node_to_slide(slide, processed_model, settings)
    
    return prs
